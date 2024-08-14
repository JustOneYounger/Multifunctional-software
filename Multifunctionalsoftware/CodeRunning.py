import os.path

import img2pdf
import matplotlib.pyplot as plt
import numpy as np
import pandas as pd
import pdf2docx
import pptx
import seaborn as sns
import win32api
import win32com.client
from MyQR import myqr
from PIL import Image
from PIL import ImageDraw
from PyQt5.QtGui import QPixmap, QImage
from PyQt5.QtWidgets import QApplication, QFrame, QHBoxLayout
from pptx.util import Inches
from pyzbar.pyzbar import decode
from spire.pdf import *

from Configs.Initconfig import *
from Configs.pdf_windows_config import *

# 支持中文
plt.rcParams['font.sans-serif'] = ['SimHei']
plt.rcParams['axes.unicode_minus'] = False


class Maincode(PdfWindowsConfig, basicconfig, initconfig):
    def __init__(self):
        basicconfig.__init__(self)
        initconfig.__init__(self)

        self.connect_buttons()
        self.newpdf_url = ""

    def connect_buttons(self):
        """
        :param:base windows
        """
        # 连接基本窗口功能
        self.pushButton_small.clicked.connect(self.becomeMin)
        self.pushButton_big.clicked.connect(self.becomeMax)
        self.pushButton_quit.clicked.connect(self.quitWindow)
        self.listWidget.clicked.connect(self.ChangeTabPage)
        self.tabWidget.currentChanged.connect(self.ChangeListItem)
        # 多窗口页面连接
        """
        :param:windows link
        """
        self.pushButton_pdf_merge.clicked.connect(self.open_mergepdf_windows)
        self.pushButton_pdf_split.clicked.connect(self.open_splitpdf_windows)
        self.pushButton_pdf_encryption.clicked.connect(self.open_encryptpdf_windows)
        # 功能连接
        """
        :param:pdf
        """
        self.pushButton_openpdf_document.clicked.connect(self.openpdf_document)
        self.pushButton_checkpdf.clicked.connect(self.checkpdf)
        self.pushButton_opensave_document.clicked.connect(self.opensave_document)
        self.pushButton_start_change_pdf.clicked.connect(self.start_change_pdf)
        self.pushButton_openwait_document.clicked.connect(self.openwait_document)
        self.pushButton_checkdocument.clicked.connect(self.checkdocument)
        self.pushButton_savepdf.clicked.connect(self.savepdf)
        self.pushButton_check_newpdf.clicked.connect(self.check_newpdf)
        self.pushButton_document_to_pdf.clicked.connect(self.document_to_pdf)
        self.pushButton_old_picture.clicked.connect(self.old_picture)
        self.pushButton_new_picture.clicked.connect(self.new_picture)
        self.pushButton_picture_type_change.clicked.connect(self.picture_type_change)
        """
        :param:qrcode
        """
        self.pushButton_qr_picture_open.clicked.connect(self.qr_picture_open)
        self.pushButton_qr_save_open.clicked.connect(self.qr_save_open)
        self.pushButton_qr_show.clicked.connect(self.myqrcode)
        self.pushButton_qr_save.clicked.connect(self.myqr_save)
        self.pushButton_upload_qr.clicked.connect(self.upload_qr)
        self.pushButton_identify_qr.clicked.connect(self.identify_qr)
        """
        :param:web
        """
        self.pushButton_insert_newweb.clicked.connect(self.insertweb)
        self.pushButton_delete_web.clicked.connect(self.deleteweb)
        """
        :param:data visual
        """
        self.pushButton_open_data_file.clicked.connect(self.open_data_file)
        self.pushButton_check_datafile.clicked.connect(self.check_datafile)
        self.comboBox_visual_types.currentIndexChanged.connect(self.visual_types_change)
        self.pushButton_visual_show.clicked.connect(self.visual_show)
        self.pushButton_visual_save_open.clicked.connect(self.visual_save_open)
        self.pushButton_visual_result_save.clicked.connect(self.visual_result_save)
        """
        :param:map
        """
        self.pushButton_map_search.clicked.connect(self.ChangeMap)

    """
    :param:soft base tab and list
    :Author:JustOneYounger
    :GitHub_url:https://github.com/JustOneYounger/Multifunctional-software
    --------------------------------------------------------------------------------------------------------------------
    """

    def ChangeTabPage(self):
        index_list = self.listWidget.currentRow()
        self.tabWidget.setCurrentIndex(index_list)

    def ChangeListItem(self):
        index_tab = self.tabWidget.currentIndex()
        if index_tab != 0:
            index_tab = index_tab
            self.listWidget.setCurrentRow(index_tab)

    """
    :param:PDF
    :Author:JustOneYounger
    :GitHub_url:https://github.com/JustOneYounger/Multifunctional-software
    --------------------------------------------------------------------------------------------------------------------
    """

    def openpdf_document(self):
        root = tk.Tk()
        root.withdraw()
        file_path = filedialog.askopenfilename(filetypes=[("PDF files", "*.pdf")])
        self.lineEdit_open_pdf.setText(file_path)

    def checkpdf(self):
        if self.lineEdit_open_pdf.text() != '':
            win32api.ShellExecute(0, "open", self.lineEdit_open_pdf.text(), "", "", 1)
        else:
            msgbox.showwarning("警告", "未选择pdf文件，无法查看")

    def opensave_document(self):
        root = tk.Tk()
        root.withdraw()
        folder_path = filedialog.askdirectory()
        self.lineEdit_save_pdf_to_document.setText(folder_path)

    def start_change_pdf(self):
        """
        :param:convert PDF to selected Word, Excel, PPT, JPG, PNG format files
        """
        if self.lineEdit_open_pdf.text() == "":
            msgbox.showwarning("警告", "pdf文件未选择")
        else:
            if self.comboBox_pdf_to_Xdocument.currentText() == "Word":

                file_path = self.lineEdit_open_pdf.text()
                file_name = os.path.basename(file_path).split('.')[0]
                save_folder = self.lineEdit_save_pdf_to_document.text()

                if save_folder != "":
                    pdf2docx.parse(file_path, save_folder + "/" + file_name + ".docx")
                    msgbox.showinfo("转换完成", "文件已保存于:%s" % (save_folder + "/" + file_name + ".docx"))
                else:
                    pdf2docx.parse(file_path, os.path.expanduser("~") + "\\Desktop\\" + file_name + ".docx")
                    msgbox.showinfo("转换完成",
                                    "文件已保存于:%s" % (os.path.expanduser("~") + "\\Desktop\\" + file_name + ".docx"))

            elif self.comboBox_pdf_to_Xdocument.currentText() == "Excel":

                file_path = self.lineEdit_open_pdf.text()
                file_name = os.path.basename(file_path).split('.')[0]
                save_folder = self.lineEdit_save_pdf_to_document.text()

                pdf = PdfDocument()
                pdf.LoadFromFile(file_path)
                convertOptions = XlsxLineLayoutOptions(True, True, False, True, False)
                pdf.ConvertOptions.SetPdfToXlsxOptions(convertOptions)

                if save_folder != "":
                    pdf.SaveToFile(save_folder + "/" + file_name + ".xlsx", FileFormat.XLSX)
                    pdf.Close()
                    msgbox.showinfo("转换完成", "文件已保存于:%s" % (save_folder + "/" + file_name + ".xlsx"))
                else:
                    pdf.SaveToFile(os.path.expanduser("~") + "\\Desktop\\" + file_name + ".xlsx", FileFormat.XLSX)
                    pdf.Close()
                    msgbox.showinfo("转换完成",
                                    "文件已保存于:%s" % (os.path.expanduser("~") + "\\Desktop\\" + file_name + ".xlsx"))

            elif self.comboBox_pdf_to_Xdocument.currentText() == "PPT":
                def change(pdf_path, ppt_name):
                    img_dir = "Images_Temp/"
                    try:
                        doc = fitz.open(pdf_path)
                        # 准备一个空的 PowerPoint 对象
                        ppt = pptx.Presentation()
                        # 遍历每一页，将其转换为图片并添加到 PowerPoint 中
                        img_paths = []
                        for page_num, page in enumerate(doc):
                            # zoom_x = 1.0  # 设置每页的水平缩放因子
                            # zoom_y = 1.0  # 设置每页的垂直缩放因子
                            # mat = fitz.Matrix(zoom_x, zoom_y)
                            pix = page.get_pixmap()
                            # 保存图片到临时目录
                            img_path = f"{img_dir}page-{page_num + 1}.png"
                            pix.save(img_path)
                            img_paths.append(img_path)
                            # 添加图片到 PowerPoint
                            layout = ppt.slide_layouts[1]  # 样式
                            slide = ppt.slides.add_slide(layout)
                            slide.shapes.add_picture(img_path, Inches(0), Inches(0), Inches(10), Inches(8))
                        ppt.save(ppt_name)
                    except Exception as e:
                        msgbox.showwarning("警告", f"转换过程中出现错误: {str(e)}")
                    finally:
                        # 关闭 PDF 文档
                        if doc:
                            doc.close()
                        # 删除临时生成的图片
                        for img_path in img_paths:
                            try:
                                os.remove(img_path)
                            except:
                                pass

                file_path = self.lineEdit_open_pdf.text()
                file_name = os.path.basename(file_path).split('.')[0]
                save_folder = self.lineEdit_save_pdf_to_document.text()

                if save_folder != "":
                    change(pdf_path=file_path, ppt_name=save_folder + "/" + file_name + ".pptx")
                    msgbox.showinfo("转换完成", "文件已保存于:%s" % (save_folder + "/" + file_name + ".pptx"))
                else:
                    change(pdf_path=file_path, ppt_name=os.path.expanduser("~") + "\\Desktop\\" + file_name + ".pptx")
                    msgbox.showinfo("转换完成",
                                    "文件已保存于:%s" % (os.path.expanduser("~") + "\\Desktop\\" + file_name + ".pptx"))

            elif self.comboBox_pdf_to_Xdocument.currentText() == "JPG":
                try:
                    # PDF文件路径
                    pdf_path = self.lineEdit_open_pdf.text()

                    # 创建输出目录
                    output_dir = self.lineEdit_save_pdf_to_document.text()
                    if output_dir == "":
                        output_dir = os.path.expanduser("~") + "\\Desktop"

                    # 将每一页图像保存为JPEG文件
                    pdf_document = fitz.open(pdf_path)
                    for page_num in range(pdf_document.page_count):
                        page = pdf_document.load_page(page_num)
                        pix = page.get_pixmap()
                        file_name = f'page_{page_num + 1}.jpg'
                        output_file = os.path.join(output_dir, file_name)
                        pix.save(output_file)
                    msgbox.showinfo("成功", "转换完成")
                except Exception as e:
                    msgbox.showwarning("警告", str(e))
            elif self.comboBox_pdf_to_Xdocument.currentText() == "PNG":
                try:
                    # PDF文件路径
                    pdf_path = self.lineEdit_open_pdf.text()

                    # 创建输出目录
                    output_dir = self.lineEdit_save_pdf_to_document.text()
                    if output_dir == "":
                        output_dir = os.path.expanduser("~") + "\\Desktop"

                    # 将每一页图像保存为JPEG文件
                    pdf_document = fitz.open(pdf_path)
                    for page_num in range(pdf_document.page_count):
                        page = pdf_document.load_page(page_num)
                        pix = page.get_pixmap()
                        file_name = f'page_{page_num + 1}.png'
                        output_file = os.path.join(output_dir, file_name)
                        pix.save(output_file)
                    msgbox.showinfo("成功", "转换完成")
                except Exception as e:
                    msgbox.showwarning("警告", str(e))

    def openwait_document(self):
        root = tk.Tk()
        root.withdraw()
        # file_path = filedialog.askopenfilename(filetypes=[
        #     ("Word Files", "*.docx;*.doc"),
        #     ("Excel Files", "*.xlsx;*.xls"),
        #     ("PowerPoint Files", "*.pptx;*.ppt"),
        #     ("JPG Files", "*.jpg;*.jpeg"),
        #     ("PNG Files", "*.png")
        # ])
        file_path = filedialog.askopenfilename()
        allow_file = ["docx", "doc", "xlsx", "xls", "pptx", "ppt", "jpg", "jpeg", "png"]
        if file_path.split('.')[-1] not in allow_file:
            msgbox.showwarning("警告", "只允许打开Word、Excel、PPT、JPG、PNG格式的文件")
            self.lineEdit_open_choose_topdf_document.clear()
        else:
            self.lineEdit_open_choose_topdf_document.setText(file_path)

    def checkdocument(self):
        if self.lineEdit_open_choose_topdf_document.text() != '':
            win32api.ShellExecute(0, "open", self.lineEdit_open_choose_topdf_document.text(), "", "", 1)
        else:
            msgbox.showwarning("警告", "未选择文件，无法查看")

    def savepdf(self):
        root = tk.Tk()
        root.withdraw()
        folder_path = filedialog.askdirectory()
        self.lineEdit_save_newpdf.setText(folder_path)

    def check_newpdf(self):
        if self.newpdf_url == "":
            msgbox.showwarning("警告", "未生成PDF文件，无法查看")
        else:
            win32api.ShellExecute(0, "open", self.newpdf_url, "", "", 1)

    def document_to_pdf(self):
        """
        :param:convert Word, Excel, PPT, JPG, PNG files to PDF files
        """
        file_path = self.lineEdit_open_choose_topdf_document.text()
        if file_path != "":
            file_name = os.path.basename(file_path).split('.')[0]
            file_type = file_path.split('.')[-1]
            save_dir = self.lineEdit_save_newpdf.text()
            if save_dir == "":
                save_dir = os.path.expanduser("~") + "\\Desktop"
            save_path = os.path.join(save_dir, file_name + ".pdf")
            self.newpdf_url = save_path

            if file_type == "docx" or file_type == "doc":
                word = win32com.client.Dispatch("Word.Application")
                doc = word.Documents.Open(file_path)
                doc.SaveAs(save_path, FileFormat=17)
                doc.Close()
                word.Quit()
                msgbox.showinfo("成功", "转换完成")
            elif file_type == "xlsx" or file_type == "xls":
                excel = win32com.client.Dispatch("Excel.Application")
                wb = excel.Workbooks.Open(file_path)
                wb.ExportAsFixedFormat(0, save_path)
                wb.Close()
                excel.Quit()
                msgbox.showinfo("成功", "转换完成")
            elif file_type == "pptx" or file_type == "ppt":
                powerpoint = win32com.client.Dispatch("PowerPoint.Application")
                ppt = powerpoint.Presentations.Open(file_path)
                ppt.SaveAs(save_path, FileFormat=32)
                ppt.Close()
                powerpoint.Quit()
                msgbox.showinfo("成功", "转换完成")
            elif file_type == "jpg" or file_type == "jpeg" or file_type == "png":
                with open(save_path, "wb") as f:
                    f.write(img2pdf.convert(file_path))
                msgbox.showinfo("成功", "转换完成")
        else:
            msgbox.showwarning("警告", "未选择文件")

    def old_picture(self):
        root = tk.Tk()
        root.withdraw()
        file_path = filedialog.askopenfilename(filetypes=[("Image files", "*.jpg;*.jpeg;*.png;*.tiff;*.bmp;*.webp")])
        self.lineEdit_old_picture.setText(file_path)

    def new_picture(self):
        root = tk.Tk()
        root.withdraw()
        folder_path = filedialog.askdirectory()
        self.lineEdit_new_picture.setText(folder_path)

    def picture_type_change(self):
        """
        :param:convert image formats to each other
        """
        picture_dict = {
            "JPEG": ".jpg",
            "PNG": ".png",
            "BMP": ".bmp",
            "TIFF": ".tiff",
            "WEBP": ".webp",
            "PPM": ".ppm",
            "PBM": ".pbm",
            "PGM": ".pgm"
        }
        new_type = self.comboBox_new_picture_type.currentText()
        file_path = self.lineEdit_old_picture.text()
        save_dir = self.lineEdit_new_picture.text()
        if save_dir == "":
            save_dir = os.path.expanduser("~") + "\\Desktop"
        if file_path == "":
            msgbox.showwarning("警告", "未选择待转换的图片")
        else:
            file_name = os.path.basename(file_path).split('.')[0]
            save_path = os.path.join(save_dir, file_name + picture_dict[new_type])
            with Image.open(file_path) as img:
                if img.mode != "RGB":
                    img = img.convert("RGB")
                img.save(save_path, format=new_type)
            msgbox.showinfo("转换完成", f"新图片已保存于{save_path}")

    """
    :param:Qrcode
    :Author:JustOneYounger
    :GitHub_url:https://github.com/JustOneYounger/Multifunctional-software
    --------------------------------------------------------------------------------------------------------------------
    """

    def qr_picture_open(self):
        root = tk.Tk()
        root.withdraw()
        file_path = filedialog.askopenfilename(filetypes=[("Image files", "*.png;*.jpg;*.bmp")])
        self.lineEdit_qr_picture.setText(file_path)

    def qr_save_open(self):
        root = tk.Tk()
        root.withdraw()
        folder_path = filedialog.askdirectory()
        self.lineEdit_qr_save.setText(folder_path)

    def myqrcode(self):
        def qr_generate_theading():
            """
            :param:generate QR code
            """
            leveldict = {
                "H(30%)": "H",
                "Q(25%)": "Q",
                "M(15%)": "M",
                "L(7%)": "L"
            }
            # 如果没有背景图
            if self.lineEdit_qr_picture.text() == "":
                myqr.run(
                    words=self.lineEdit_qr_word.text(),
                    level=leveldict[self.comboBox_qr_level.currentText()],
                    save_name='qrcode_temp.png',
                    save_dir='./Temp/'
                )
            else:
                myqr.run(
                    words=self.lineEdit_qr_word.text(),
                    level=leveldict[self.comboBox_qr_level.currentText()],
                    picture=self.lineEdit_qr_picture.text(),
                    colorized=True,
                    save_name='qrcode_temp.png',
                    save_dir='./Temp/'
                )

        def qr_show_threading():
            """
            :param:load image into container
            """
            # 加载图片到 QPixmap 对象
            pixmap = QPixmap('./Temp/qrcode_temp.png')
            # 设置 QLabel 的 pixmap
            self.label_showmyqr.setPixmap(pixmap)
            self.label_showmyqr.setScaledContents(True)
            file_size = os.path.getsize('./Temp/qrcode_temp.png')
            if file_size < 1024:
                self.label_qrcode_picture_size.setText(str(file_size) + " B")
            elif file_size >= 1024:
                file_size = file_size / 1024
                self.label_qrcode_picture_size.setText(str(round(file_size, 3)) + " KB")
            elif file_size >= 1024 * 1024:
                file_size = file_size / (1024 * 1024)
                self.label_qrcode_picture_size.setText(str(round(file_size, 3)) + " MB")

        generate_thead = threading.Thread(target=qr_generate_theading)
        show_thread = threading.Thread(target=qr_show_threading)
        generate_thead.start()
        generate_thead.join()
        show_thread.start()

    def myqr_save(self):
        """
        :param:save the generated QR code
        """
        if self.lineEdit_qr_save.text() != "":
            if os.path.exists("./Temp/qrcode_temp.png"):
                filename = os.path.basename("./Temp/qrcode_temp.png")
                save_path = os.path.join(self.lineEdit_qr_save.text(), filename)
                shutil.copyfile("./Temp/qrcode_temp.png", save_path)
                msgbox.showinfo("保存成功", f"二维码已保存于{self.lineEdit_qr_save.text() + '/' + filename}")
            else:
                msgbox.showwarning(title="警告", message="二维码还未生成，无法保存")
        else:
            if os.path.exists("./Temp/qrcode_temp.png"):
                filename = os.path.basename("./Temp/qrcode_temp.png")
                save_path = os.path.join(os.path.expanduser("~") + "/Desktop", filename)
                shutil.copyfile("./Temp/qrcode_temp.png", save_path)
                msgbox.showinfo("保存成功", "二维码已保存于桌面")
            else:
                msgbox.showwarning(title="警告", message="二维码还未生成，无法保存")

    def upload_qr(self):
        """
        :param:upload an image containing a QR code
        """
        root = tk.Tk()
        root.withdraw()
        file_path = filedialog.askopenfilename(filetypes=[("Image files", "*.png;*.jpg;*.bmp")])
        if file_path != '':
            self.upload_qr_file_path = file_path
            self.textBrowser_identify_result.clear()
            # 打开图片文件
            img = Image.open(file_path)
            # 计算宽高比例，保持宽度与高度比例不变
            width_percent = (250 / float(img.size[1]))
            new_width = int((float(img.size[0]) * float(width_percent)))
            # 调整图片大小
            resized_img = img.resize((new_width, 250), Image.LANCZOS)
            # 转换图片数据并装配label
            image = resized_img.convert("RGBA")
            image_data = image.tobytes("raw", "RGBA")
            q_image = QImage(image_data, image.size[0], image.size[1], QImage.Format_RGBA8888)
            pixmap = QPixmap.fromImage(q_image)
            self.label_upload_qrpicture.setPixmap(pixmap)
            self.label_upload_qrpicture.setScaledContents(True)

    def identify_qr(self):
        def draw_picture_threading():
            """
            :param:decode
            """
            img = Image.open(self.upload_qr_file_path)
            if img.mode != 'RGB':
                img = img.convert("RGB")
            draw = ImageDraw.Draw(img)
            decode_data = decode(img)
            for barcode in decode_data:
                self.textBrowser_identify_result.append(f'类型为:\n{barcode.type}')
                url_string = barcode.data.decode('utf-8')
                self.textBrowser_identify_result.append(f'解码信息为:\n{url_string}')
                self.textBrowser_identify_result.append('')
                pts = np.array(barcode.polygon, np.int32)
                draw.polygon([tuple(pt) for pt in pts], outline='red', width=5)
            img.save('./Temp/upload_draw_qrcode.png')

        def show_draw_threading():
            # 打开图片文件
            img = Image.open('./Temp/upload_draw_qrcode.png')
            # 计算宽高比例，保持宽度与高度比例不变
            width_percent = (250 / float(img.size[1]))
            new_width = int((float(img.size[0]) * float(width_percent)))
            # 调整图片大小
            resized_img = img.resize((new_width, 250), Image.LANCZOS)
            # 转换图片数据并装配label
            image = resized_img.convert("RGBA")
            image_data = image.tobytes("raw", "RGBA")
            q_image = QImage(image_data, image.size[0], image.size[1], QImage.Format_RGBA8888)
            pixmap = QPixmap.fromImage(q_image)
            self.label_upload_qrpicture.setPixmap(pixmap)
            self.label_upload_qrpicture.setScaledContents(True)

        draw_picture_thread = threading.Thread(target=draw_picture_threading)
        show_draw_thread = threading.Thread(target=show_draw_threading)
        draw_picture_thread.start()
        draw_picture_thread.join()
        show_draw_thread.start()

    """
    :param:Web
    :Author:JustOneYounger
    :GitHub_url:https://github.com/JustOneYounger/Multifunctional-software
    --------------------------------------------------------------------------------------------------------------------
    """

    def insertweb(self):
        """
        :param:insert a new web into tabwidget
        """
        url = self.lineEdit_web_url_input.text()
        if url:
            max_index = self.tabWidget_webs.count()

            new_web = WebEngineView()

            frame_web = QFrame()
            frame_layout = QVBoxLayout()
            frame_layout.setContentsMargins(0, 0, 0, 0)
            frame_layout.addWidget(new_web)
            frame_web.setLayout(frame_layout)

            tab_layout = QHBoxLayout()
            tab_layout.setContentsMargins(0, 0, 0, 0)
            self.tabWidget_webs.setLayout(tab_layout)
            self.tabWidget_webs.addTab(frame_web, f'{max_index}')

            new_web.load(QUrl(url))

            self.tabWidget_webs.setCurrentIndex(max_index)
            self.spinBox_web_index.setValue(max_index)

    def deleteweb(self):
        """
        :param:delete the selected page
        """
        max_index = self.tabWidget_webs.count()
        current_index = self.tabWidget_webs.currentIndex()
        if current_index >= 0:
            self.tabWidget_webs.removeTab(current_index)
            for i in range(current_index, max_index):
                self.tabWidget_webs.setTabText(i, f'{i}')

    """
    :param:Data Visual
    :Author:JustOneYounger
    :GitHub_url:https://github.com/JustOneYounger/Multifunctional-software
    --------------------------------------------------------------------------------------------------------------------
    """

    def open_data_file(self):
        root = tk.Tk()
        root.withdraw()
        file_path = filedialog.askopenfilename(filetypes=[('Excel files', '*.xls;*.xlsx'), ('CSV files', '*.csv')])
        self.lineEdit_datafile_url.setText(file_path)

    def check_datafile(self):
        if self.lineEdit_datafile_url.text() == "":
            msgbox.showwarning("警告", "未选择数据文件，无法打开查看")
        else:
            win32api.ShellExecute(0, "open", self.lineEdit_datafile_url.text(), "", "", 1)

    def visual_types_change(self):
        if self.comboBox_visual_types.currentText() == '折线图':
            self.doubleSpinBox_plot_linewidth.setEnabled(True)
            self.doubleSpinBox_bar_width.setEnabled(False)
            self.comboBox_hot_color.setEnabled(False)
            self.checkBox_heatmap_annot.setEnabled(False)
        elif self.comboBox_visual_types.currentText() == '纵向条形图' or self.comboBox_visual_types.currentText() == '横向条形图':
            self.doubleSpinBox_plot_linewidth.setEnabled(False)
            self.doubleSpinBox_bar_width.setEnabled(True)
            self.comboBox_hot_color.setEnabled(False)
            self.checkBox_heatmap_annot.setEnabled(False)
        elif self.comboBox_visual_types.currentText() == '热力图':
            self.doubleSpinBox_plot_linewidth.setEnabled(False)
            self.doubleSpinBox_bar_width.setEnabled(False)
            self.comboBox_hot_color.setEnabled(True)
            self.checkBox_heatmap_annot.setEnabled(True)
        else:
            self.doubleSpinBox_plot_linewidth.setEnabled(False)
            self.doubleSpinBox_bar_width.setEnabled(False)
            self.comboBox_hot_color.setEnabled(False)
            self.checkBox_heatmap_annot.setEnabled(False)

    def visual_show(self):
        def draw_visual_threading():
            heatmap_cmp_dict = {
                "默认": "viridis",
                "浅蓝-深蓝": "Blues",
                "浅绿-深绿": "Greens",
                "浅橘-深橘": "Oranges",
                "浅红-深红": "Reds",
                "浅紫-深紫": "Purples",
                "黄-橘-棕": "YlOrBr",
                "黄绿-深蓝": "YlGnBu",
                "红-黄-蓝": "RdYlBu",
                "彩虹色": "Spectral",
                "黄-黑-红": "inferno",
                "粉红-黄-紫": "plasma",
                "紫-黑-红": "magma",
                "绿-黄": "cividis"
            }
            if self.lineEdit_datafile_url.text() != '':
                if self.lineEdit_datafile_url.text().split('.')[-1] == 'csv':
                    if self.checkBox_read_headline.isChecked() == False:
                        data = pd.read_csv(self.lineEdit_datafile_url.text(), header=None)
                    else:
                        data = pd.read_csv(self.lineEdit_datafile_url.text())
                        labels = list(data.keys())
                    data = np.array(data)
                elif self.lineEdit_datafile_url.text().split('.')[-1] == 'xlsx' or \
                        self.lineEdit_datafile_url.text().split('.')[-1] == 'xls':
                    if self.checkBox_read_headline.isChecked() == False:
                        data = pd.read_excel(self.lineEdit_datafile_url.text(), header=None)
                    else:
                        data = pd.read_excel(self.lineEdit_datafile_url.text())
                        labels = list(data.keys())
                    data = np.array(data)
                fig = plt.figure()

                if self.lineEdit_title.text() != '':
                    plt.title(self.lineEdit_title.text())
                if self.lineEdit_X_name.text() != '':
                    plt.xlabel(self.lineEdit_X_name.text())
                if self.lineEdit_Y_name.text() != '':
                    plt.ylabel(self.lineEdit_Y_name.text())
                plt.grid(self.checkBox_grid_start.isChecked())

                if self.comboBox_visual_types.currentText() == '折线图':
                    try:
                        if data.shape[0] == 1:
                            plt.plot(
                                data[0],
                                linewidth=self.doubleSpinBox_plot_linewidth.value()
                            )
                        else:
                            plt.plot(
                                data,
                                linewidth=self.doubleSpinBox_plot_linewidth.value()
                            )
                        plt.savefig("Temp/visual.png")
                    except Exception as e:
                        self.textBrowser_visual_error.append(str(e) + '\n' + "请检查相关设置")
                        self.textBrowser_visual_error.append("\n如遇困难，可联系客服19550112627")
                elif self.comboBox_visual_types.currentText() == '纵向条形图':
                    if data.shape[0] == 1:
                        width = self.doubleSpinBox_bar_width.value()
                        for i in range(data.shape[1]):
                            plt.bar(
                                labels[i],
                                data[:, i],
                                width=width
                            )
                        plt.savefig("Temp/visual.png")
                    else:
                        self.textBrowser_visual_error.append(
                            "目前只接受一维数据，请修改数据源为首行标签行且第二行为数据")
                        self.textBrowser_visual_error.append("\n如遇困难，可联系客服19550112627")
                elif self.comboBox_visual_types.currentText() == '横向条形图':
                    if data.shape[0] == 1:
                        width = self.doubleSpinBox_bar_width.value()
                        for i in range(data.shape[1]):
                            plt.barh(
                                labels[i],
                                data[:, i],
                                height=width
                            )
                        plt.savefig("Temp/visual.png")
                    else:
                        self.textBrowser_visual_error.append(
                            "目前只接受一维数据，请修改数据源为首行标签行且第二行为数据")
                        self.textBrowser_visual_error.append("\n如遇困难，可联系客服19550112627")
                elif self.comboBox_visual_types.currentText() == '饼图':
                    if data.shape[0] == 1:
                        plt.pie(
                            data[0],
                            labels=labels,
                            autopct='%.2f%%'
                        )
                        plt.savefig("Temp/visual.png")
                    else:
                        self.textBrowser_visual_error.append(
                            "饼图只能接受一维数据，请修改数据源为首行标签行且第二行为数据")
                        self.textBrowser_visual_error.append("\n如遇困难，可联系客服19550112627")
                elif self.comboBox_visual_types.currentText() == '热力图':
                    try:
                        sns.heatmap(
                            data,
                            cmap=heatmap_cmp_dict[self.comboBox_hot_color.currentText()],
                            annot=self.checkBox_heatmap_annot.isChecked()
                        )
                        plt.savefig("Temp/visual.png")
                    except Exception as e:
                        self.textBrowser_visual_error.append(str(e) + '\n' + "请检查相关设置")
                        self.textBrowser_visual_error.append("\n如遇困难，可联系客服19550112627")

        def show_visual_threading():
            if os.path.exists('./Temp/visual.png'):
                pixmap = QPixmap('./Temp/visual.png')
                self.label_visual_show.setPixmap(pixmap)
                self.label_visual_show.setScaledContents(True)
                file_size = os.path.getsize('./Temp/visual.png')
                if file_size < 1024:
                    self.label_visual_size.setText("大小为： " + str(file_size) + " B")
                elif file_size >= 1024:
                    file_size = file_size / 1024
                    self.label_visual_size.setText("大小为： " + str(round(file_size, 3)) + " KB")
                elif file_size >= 1024 * 1024:
                    file_size = file_size / (1024 * 1024)
                    self.label_visual_size.setText("大小为： " + str(round(file_size, 3)) + " MB")
            else:
                self.textBrowser_visual_error.append("运行发生错误，请检查相关设置")

        self.textBrowser_visual_error.setText(None)
        draw_visual_thread = threading.Thread(target=draw_visual_threading)
        show_visual_thread = threading.Thread(target=show_visual_threading)
        draw_visual_thread.start()
        draw_visual_thread.join()
        show_visual_thread.start()

    def visual_save_open(self):
        root = tk.Tk()
        root.withdraw()
        folder_path = filedialog.askdirectory()
        self.lineEdit_save_visual_url.setText(folder_path)

    def visual_result_save(self):
        if self.lineEdit_save_visual_url.text() != "":
            if os.path.exists("./Temp/visual.png"):
                filename = os.path.basename("./Temp/visual.png")
                save_path = os.path.join(self.lineEdit_save_visual_url.text(), filename)
                shutil.copyfile("./Temp/visual.png", save_path)
                msgbox.showinfo("保存成功", f"图像已保存于{self.lineEdit_save_visual_url.text() + '/' + filename}")
            else:
                msgbox.showwarning(title="警告", message="图像还未生成，无法保存")
        else:
            if os.path.exists("./Temp/visual.png"):
                filename = os.path.basename("./Temp/visual.png")
                save_path = os.path.join(os.path.expanduser("~") + "/Desktop", filename)
                shutil.copyfile("./Temp/visual.png", save_path)
                msgbox.showinfo("保存成功", "图像已保存于桌面")
            else:
                msgbox.showwarning(title="警告", message="图像还未生成，无法保存")

    """
    :param:Map
    :Author:JustOneYounger
    :GitHub_url:https://github.com/JustOneYounger/Multifunctional-software
    --------------------------------------------------------------------------------------------------------------------
    """

    def ChangeMap(self):
        """
        :param:update the map based on parameter selection
        """
        Map = folium.Map([float(self.lineEdit_latitude.text()), float(self.lineEdit_longitude.text())],
                         tiles=self.mymapdict[self.comboBox_map_types.currentText()],
                         attr=self.comboBox_map_types.currentText(),
                         zoom_start=self.spinBox_map_proportion.value(),
                         )
        Map.add_child(folium.LatLngPopup())  # 显示鼠标点击点经纬度
        Map_html = Map.get_root().render()  # Map转化为html
        self.qwebengine_map.setHtml(Map_html)
        self.textBrowser_map.append("---------------------------------------")
        self.textBrowser_map.append(f"当前地图类型:{self.comboBox_map_types.currentText()}")
        self.textBrowser_map.append(f"当前经度为:{self.lineEdit_longitude.text()}")
        self.textBrowser_map.append(f"当前纬度为:{self.lineEdit_latitude.text()}")
        self.textBrowser_map.append(f"当前缩放比为:{self.spinBox_map_proportion.value()}")


if __name__ == "__main__":
    app = QApplication(sys.argv)
    mainrun = Maincode()
    mainrun.show()
    sys.exit(app.exec_())
